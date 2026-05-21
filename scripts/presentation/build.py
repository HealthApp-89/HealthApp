"""
Build the Apex Health OS UX presentation.

Pipeline:
  1. Render mockup HTML scenes (scenes/) via headless chromium → shots/.
  2. Pull real authenticated app screenshots from shots-real/ (produced
     separately by capture-real.mjs).
  3. Assemble a .pptx with cover, "why this exists", executive summary,
     walkthrough, and mental-model slides at docs/Apex-Health-OS-Walkthrough.pptx.

Usage:
  node --env-file=.env.local scripts/presentation/capture-real.mjs   # once
  python3 scripts/presentation/build.py                              # rebuild deck
"""
from __future__ import annotations

import subprocess
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from PIL import Image

ROOT = Path(__file__).resolve().parents[2]
SCENES = Path(__file__).parent / "scenes"
SHOTS = Path(__file__).parent / "shots"            # mockups
REAL = Path(__file__).parent / "shots-real"        # real app captures
OUT = ROOT / "docs" / "Apex-Health-OS-Walkthrough.pptx"
OUT.parent.mkdir(parents=True, exist_ok=True)
SHOTS.mkdir(parents=True, exist_ok=True)

CHROME = (
    "/Users/abdelouahedelbied/Library/Caches/ms-playwright/"
    "chromium-1217/chrome-mac-arm64/Google Chrome for Testing.app/"
    "Contents/MacOS/Google Chrome for Testing"
)

# Mockup scenes (HTML → PNG). Only the ones needed by the current deck.
MOCKUP_VIEWPORTS = {
    "onboarding":   (470, 920),
    "integrations": (470, 920),
    "intake":       (470, 980),
    "brief":        (470, 1280),
    "weekly":       (470, 1060),
}


def screenshot_mockup(scene: str, vw: int, vh: int) -> Path:
    html = SCENES / f"{scene}.html"
    out = SHOTS / f"{scene}.png"
    if not html.exists():
        raise FileNotFoundError(html)
    subprocess.run(
        [
            CHROME,
            "--headless=new",
            "--disable-gpu",
            "--hide-scrollbars",
            "--no-sandbox",
            "--force-device-scale-factor=2",
            f"--window-size={vw},{vh}",
            f"--screenshot={out}",
            f"file://{html}",
        ],
        check=True,
        capture_output=True,
    )
    im = Image.open(out).convert("RGBA")
    bbox = im.getbbox()
    if bbox:
        im.crop(bbox).save(out)
    return out


def render_all_mockups() -> dict[str, Path]:
    shots: dict[str, Path] = {}
    for scene, (vw, vh) in MOCKUP_VIEWPORTS.items():
        print(f"  rendering mockup {scene}…")
        shots[scene] = screenshot_mockup(scene, vw, vh)
    return shots


# ─── pptx assembly ────────────────────────────────────────────────────────────

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Brand palette (mirrors lib/ui/theme.ts)
BG          = (241, 242, 246)
SURFACE     = (255, 255, 255)
STRONG      = (15, 20, 48)
MID         = (74, 77, 98)
MUTED       = (122, 126, 149)
FAINT       = (144, 148, 168)
ACCENT      = (79, 93, 255)
ACCENT_SOFT = (231, 234, 255)
ACCENT_DEEP = (58, 71, 232)
SUCCESS     = (20, 184, 112)
WARNING     = (245, 158, 11)
DANGER      = (239, 68, 68)


def rgb(t):
    return RGBColor(*t)


def set_slide_bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb(color)


def add_text(slide, x, y, w, h, text, *,
             size=14, bold=False, color=STRONG, align=PP_ALIGN.LEFT,
             font="Calibri", italic=False, line_spacing=None):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    para = tf.paragraphs[0]
    para.alignment = align
    if line_spacing:
        para.line_spacing = line_spacing
    run = para.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = rgb(color)
    return tb


def add_multitext(slide, x, y, w, h, lines, *, line_spacing=None):
    """lines: list of (text, dict). dict supports size/bold/italic/color/align."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    for i, (text, opts) in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = opts.get("align", PP_ALIGN.LEFT)
        if line_spacing:
            para.line_spacing = line_spacing
        if "space_after" in opts:
            para.space_after = Pt(opts["space_after"])
        run = para.add_run()
        run.text = text
        run.font.name = "Calibri"
        run.font.size = Pt(opts.get("size", 14))
        run.font.bold = opts.get("bold", False)
        run.font.italic = opts.get("italic", False)
        run.font.color.rgb = rgb(opts.get("color", STRONG))
    return tb


def add_rect(slide, x, y, w, h, *, fill=None, line=None, radius=False, line_width=0.75):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        x, y, w, h,
    )
    if radius:
        shape.adjustments[0] = radius if isinstance(radius, float) else 0.08
    if fill is not None:
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb(fill)
    else:
        shape.fill.background()
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = rgb(line)
        shape.line.width = Pt(line_width)
    return shape


def add_pill(slide, x, y, w, h, text, *, fill=ACCENT, text_color=(255, 255, 255), size=10):
    p = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    p.adjustments[0] = 0.5
    p.fill.solid()
    p.fill.fore_color.rgb = rgb(fill)
    p.line.fill.background()
    tf = p.text_frame
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.CENTER
    run = para.add_run()
    run.text = text
    run.font.name = "Calibri"
    run.font.size = Pt(size)
    run.font.bold = True
    run.font.color.rgb = rgb(text_color)


def add_bullets(slide, x, y, w, h, items, *, size=14, color=MID, gap=4, glyph_color=ACCENT):
    """items: list of (head_bold, rest) or plain str."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    for i, item in enumerate(items):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = PP_ALIGN.LEFT
        para.space_after = Pt(gap)
        glyph = para.add_run()
        glyph.text = "●   "
        glyph.font.name = "Calibri"
        glyph.font.size = Pt(size)
        glyph.font.color.rgb = rgb(glyph_color)
        glyph.font.bold = True
        if isinstance(item, tuple):
            head, tail = item
            r1 = para.add_run()
            r1.text = head
            r1.font.name = "Calibri"
            r1.font.size = Pt(size)
            r1.font.bold = True
            r1.font.color.rgb = rgb(STRONG)
            r2 = para.add_run()
            r2.text = tail
            r2.font.name = "Calibri"
            r2.font.size = Pt(size)
            r2.font.color.rgb = rgb(color)
        else:
            r = para.add_run()
            r.text = item
            r.font.name = "Calibri"
            r.font.size = Pt(size)
            r.font.color.rgb = rgb(color)
    return tb


def add_image_fit(slide, path: Path, x, y, max_w, max_h):
    """Place image preserving aspect ratio, fitting within (max_w, max_h)."""
    with Image.open(path) as im:
        iw, ih = im.size
    target_w = max_w
    target_h = int(max_w * ih / iw)
    if target_h > max_h:
        target_h = max_h
        target_w = int(max_h * iw / ih)
    cx = x + (max_w - target_w) // 2
    cy = y + (max_h - target_h) // 2
    slide.shapes.add_picture(str(path), cx, cy, width=target_w, height=target_h)


def add_phone_frame(slide, path: Path, x, y, max_w, max_h):
    """Wrap a raw app screenshot in a rounded phone-like border."""
    # Frame slightly larger than image
    pad = Inches(0.04)
    add_rect(
        slide,
        x - pad, y - pad,
        max_w + pad * 2, max_h + pad * 2,
        fill=STRONG, radius=0.04,
    )
    add_image_fit(slide, path, x, y, max_w, max_h)


# ─── chrome ───────────────────────────────────────────────────────────────────

def add_brand_rail(slide):
    """Thin indigo rail down the left edge — present on every content slide."""
    rail = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.15), SLIDE_H)
    rail.fill.solid()
    rail.fill.fore_color.rgb = rgb(ACCENT)
    rail.line.fill.background()


def add_top_band(slide):
    """Faint indigo band along the top — adds brand presence without noise."""
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.15), Inches(0), SLIDE_W - Inches(0.15), Inches(0.12))
    band.fill.solid()
    band.fill.fore_color.rgb = rgb(ACCENT_SOFT)
    band.line.fill.background()


def add_page_chip(slide, text: str):
    w, h = Inches(0.95), Inches(0.32)
    x, y = SLIDE_W - w - Inches(0.5), Inches(0.35)
    add_pill(slide, x, y, w, h, text, fill=ACCENT, text_color=(255, 255, 255), size=9)


def add_slide_header(slide, eyebrow: str, title: str):
    add_text(
        slide, Inches(0.6), Inches(0.45), Inches(8.5), Inches(0.32),
        eyebrow.upper(), size=11, bold=True, color=ACCENT,
    )
    add_text(
        slide, Inches(0.6), Inches(0.78), Inches(11.5), Inches(0.85),
        title, size=28, bold=True, color=STRONG,
    )


def add_footer(slide, text="APEX HEALTH OS — User experience walkthrough"):
    add_text(
        slide, Inches(0.6), SLIDE_H - Inches(0.42),
        Inches(12.0), Inches(0.3),
        text, size=9, color=FAINT,
    )


def content_chrome(slide, page_num: int, total: int):
    set_slide_bg(slide)
    add_brand_rail(slide)
    add_top_band(slide)
    add_page_chip(slide, f"{page_num} / {total}")
    add_footer(slide)


# ─── slide builders ───────────────────────────────────────────────────────────

def slide_title(prs, hero: Path):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s, color=STRONG)
    # Indigo rail extension — left third is solid accent for contrast
    add_rect(s, Inches(0), Inches(0), Inches(5.2), SLIDE_H, fill=ACCENT)
    add_rect(s, Inches(5.0), Inches(0), Inches(0.2), SLIDE_H, fill=ACCENT_DEEP)

    add_text(
        s, Inches(0.7), Inches(0.7), Inches(4.0), Inches(0.4),
        "APEX HEALTH OS", size=12, bold=True, color=(231, 234, 255),
    )
    add_multitext(s, Inches(0.7), Inches(2.4), Inches(4.5), Inches(3.2), [
        ("User", {"size": 56, "bold": True, "color": (255, 255, 255)}),
        ("experience", {"size": 56, "bold": True, "color": (255, 255, 255)}),
        ("walkthrough.", {"size": 56, "bold": True, "color": ACCENT_SOFT}),
    ], line_spacing=1.0)
    add_text(
        s, Inches(0.7), Inches(6.3), Inches(4.4), Inches(0.5),
        "From first open to nightly close-out.",
        size=15, color=(220, 224, 255),
    )
    add_text(
        s, Inches(0.7), Inches(6.7), Inches(4.4), Inches(0.4),
        f"v1.1 · May 16, 2026",
        size=10, color=(180, 188, 220),
    )

    # Right side: hero brief
    add_text(
        s, Inches(5.7), Inches(0.85), Inches(7.0), Inches(0.45),
        "ONE COACH · FIVE STREAMS · THREE RHYTHMS",
        size=11, bold=True, color=(180, 188, 220),
    )
    add_text(
        s, Inches(5.7), Inches(1.3), Inches(7.0), Inches(1.5),
        "A 5-minute tour\nof the daily ritual.",
        size=26, bold=True, color=(255, 255, 255), line_spacing=1.1,
    )
    add_image_fit(s, hero, Inches(6.0), Inches(2.6), Inches(6.5), Inches(4.7))


def slide_why(prs, page_num: int, total: int):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    content_chrome(s, page_num, total)
    add_slide_header(s, "Why this exists",
                     "Generic trackers measure. Generic coaches advise.\nNeither knows you.")

    # Left: the problem — three pain points in cards
    pains = [
        ("Data scattered across 5 apps.",
         "WHOOP shows recovery. Withings shows weight. Strong shows lifts. Yazio shows macros. None of them talk to each other — so no app can answer the only question that matters: am I on track?"),
        ("Generic coaches don't know you.",
         "A LLM with no context tells everyone the same things. No medical history, no GLP-1 phase, no injury list, no “the kids have a play tonight so I'm skipping the gym” — so the advice is generic and ignorable."),
        ("Plans drift. Sundays do not.",
         "Without a Sunday plan and a Monday morning brief, the week becomes whatever Tuesday felt like. Adherence dies in the dark — most apps don't even ask if you did what you said you would."),
    ]
    y = Inches(2.05)
    for head, body in pains:
        card = add_rect(s, Inches(0.55), y, Inches(7.4), Inches(1.45), fill=SURFACE, line=(232, 234, 243), radius=0.04, line_width=0.5)
        card.shadow.inherit = False
        # Small accent square
        add_rect(s, Inches(0.55), y, Inches(0.08), Inches(1.45), fill=ACCENT)
        add_text(s, Inches(0.85), y + Inches(0.15), Inches(7.0), Inches(0.4),
                 head, size=16, bold=True, color=STRONG)
        add_text(s, Inches(0.85), y + Inches(0.55), Inches(7.0), Inches(0.85),
                 body, size=11, color=MID, line_spacing=1.25)
        y += Inches(1.6)

    # Right column: thesis box
    thesis_x = Inches(8.4)
    thesis_w = Inches(4.5)
    thesis_y = Inches(2.05)
    thesis_h = Inches(4.95)
    box = add_rect(s, thesis_x, thesis_y, thesis_w, thesis_h, fill=ACCENT, radius=0.04)
    box.shadow.inherit = False
    add_text(s, thesis_x + Inches(0.3), thesis_y + Inches(0.35), thesis_w - Inches(0.6), Inches(0.35),
             "THE THESIS", size=11, bold=True, color=(220, 224, 255))
    add_text(s, thesis_x + Inches(0.3), thesis_y + Inches(0.75), thesis_w - Inches(0.6), Inches(2.2),
             "One athlete.\nOne coach.\nOne file.",
             size=30, bold=True, color=(255, 255, 255), line_spacing=1.15)
    add_text(s, thesis_x + Inches(0.3), thesis_y + Inches(2.8), thesis_w - Inches(0.6), Inches(2.1),
             "Apex Health OS is a single-user app for a single user. It owns the data pipes, holds the athlete profile as a durable client file, and runs a server-side coach that has read that file before saying a single word.",
             size=12, color=(231, 234, 255), line_spacing=1.4)


def slide_summary(prs, page_num: int, total: int):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    content_chrome(s, page_num, total)
    add_slide_header(s, "Executive summary", "Six commitments shape the product.")

    items = [
        ("Single-user by design.",
         "Built for one athlete with one profile. No multi-tenant tax, no cohort assumptions, no compromises in the prompts."),
        ("Five sources, one source-of-truth.",
         "WHOOP · Withings · Apple Health (Garmin) · Strong · Yazio. Each metric has exactly one owner enforced in code — no overwrites."),
        ("Acknowledged profile as durable client file.",
         "Onboarding produces an immutable v1 covering health, training, lifestyle, nutrition, sleep, and the goal-with-why. Revisions create v2."),
        ("Plan generated, not improvised.",
         "Deterministic composer produces strength + nutrition + sleep + recovery. The AI narrates; it never fabricates the prescription."),
        ("Three rhythms, one document, one chat.",
         "Morning brief sets the day. The proactive layer pings when a trigger fires (plateau · off-pace · HRV). The Sunday review document closes the week."),
        ("Sundays own the week. Blocks own the cycle.",
         "Versioned Weekly Review with recap → reconfirm → trends → prescription, mid-week swaps with adherence math, monthly body measurements, 5-week mesocycle blocks."),
    ]
    # 2-column grid
    col_w = Inches(5.85)
    col_h = Inches(1.45)
    positions = [
        (Inches(0.55), Inches(1.95)),
        (Inches(6.85), Inches(1.95)),
        (Inches(0.55), Inches(3.55)),
        (Inches(6.85), Inches(3.55)),
        (Inches(0.55), Inches(5.15)),
        (Inches(6.85), Inches(5.15)),
    ]
    for (x, y), (head, body) in zip(positions, items):
        card = add_rect(s, x, y, col_w, col_h, fill=SURFACE, line=(232, 234, 243), radius=0.04, line_width=0.5)
        card.shadow.inherit = False
        # Number badge
        n = positions.index((x, y)) + 1
        add_pill(s, x + Inches(0.25), y + Inches(0.22), Inches(0.42), Inches(0.42),
                 str(n), fill=ACCENT, text_color=(255, 255, 255), size=14)
        add_text(s, x + Inches(0.85), y + Inches(0.22), col_w - Inches(1.1), Inches(0.4),
                 head, size=15, bold=True, color=STRONG)
        add_text(s, x + Inches(0.85), y + Inches(0.62), col_w - Inches(1.1), Inches(0.8),
                 body, size=11, color=MID, line_spacing=1.3)


def slide_first_open(prs, page_num: int, total: int, onboarding: Path):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    content_chrome(s, page_num, total)
    add_slide_header(s, "Step 1 · First open",
                     "Register, build the athlete profile, set the goal.")
    add_bullets(s, Inches(0.6), Inches(2.0), Inches(7.0), Inches(5.0), [
        ("Sign up. ", "Email + password (Supabase Auth). Single-user app — no team setup."),
        ("6-step onboarding wizard. ", "Health → Training → Lifestyle → Nutrition → Sleep → Goals."),
        ("Acknowledge to lock the profile. ", "v1 is immutable; future edits create v2 with the prior version superseded."),
        ("Goal + why. ", "The narrative is what gets surfaced when motivation dips. Why beats what."),
        ("Then: a 5-beat coach intake. ", "Deterministic plan composer produces strength template, macros, sleep window, recovery."),
        ("Block setup. ", "5-week mesocycle goal (e.g. “+5 kg squat 1RM, hold 80 kg LBM”)."),
    ], size=13, gap=6)
    add_image_fit(s, onboarding, Inches(8.4), Inches(1.85), Inches(4.5), Inches(5.3))


def slide_connect(prs, page_num: int, total: int, profile_real: Path, integrations: Path):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    content_chrome(s, page_num, total)
    add_slide_header(s, "Step 2 · Connect your streams",
                     "Five sources, one daily-log row per day.")
    add_bullets(s, Inches(0.6), Inches(2.0), Inches(5.8), Inches(5.0), [
        ("WHOOP. ", "OAuth → recovery, sleep, strain, HRV, RHR, SpO₂."),
        ("Withings. ", "OAuth → weight, body fat, muscle mass, hydration."),
        ("Apple Health (Garmin). ", "Per-user bearer token + Shortcut → steps, calories, distance."),
        ("Strong. ", "Webhook ingests every set, every PR."),
        ("Yazio. ", "Webhook ingests meals → macros bar in real time."),
        ("Owners in code, not docs. ", "Each metric has one source-of-truth integration; the merge layer enforces it."),
    ], size=13, gap=6)
    # Two phones side-by-side
    add_text(s, Inches(7.0), Inches(1.85), Inches(2.4), Inches(0.3),
             "PROFILE (REAL)", size=10, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_image_fit(s, profile_real, Inches(6.7), Inches(2.2), Inches(2.9), Inches(5.0))
    add_text(s, Inches(9.9), Inches(1.85), Inches(2.4), Inches(0.3),
             "INTEGRATIONS", size=10, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_image_fit(s, integrations, Inches(9.7), Inches(2.2), Inches(3.0), Inches(5.0))


def slide_morning(prs, page_num: int, total: int, intake: Path, brief: Path):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    content_chrome(s, page_num, total)
    add_slide_header(s, "Morning ritual",
                     "Two screens, two minutes — the day's plan is set.")

    # Two phones with labels
    add_text(s, Inches(0.6), Inches(1.9), Inches(4.0), Inches(0.3),
             "INTAKE", size=10, bold=True, color=ACCENT)
    add_image_fit(s, intake, Inches(0.6), Inches(2.2), Inches(4.0), Inches(5.0))

    add_text(s, Inches(5.0), Inches(1.9), Inches(4.5), Inches(0.3),
             "MORNING BRIEF (hero)", size=10, bold=True, color=ACCENT)
    add_image_fit(s, brief, Inches(5.0), Inches(2.2), Inches(4.5), Inches(5.0))

    # Right column: callouts
    add_text(s, Inches(9.9), Inches(1.9), Inches(3.0), Inches(0.3),
             "WHAT'S INSIDE", size=10, bold=True, color=ACCENT)
    add_bullets(s, Inches(9.9), Inches(2.25), Inches(3.2), Inches(4.9), [
        ("Yesterday recap. ", "Sleep, strain, steps, macros hit/missed."),
        ("Readiness band. ", "High / moderate / low + top drivers."),
        ("Today's session. ", "Or recovery focus on rest days."),
        ("Macros target. ", "Adapts to plan: GLP-1, cut, recomp."),
        ("Coach advice. ", "One paragraph, adaptive to today's flags."),
        ("Tonight's sleep target. ", "Lights-out time."),
    ], size=11, gap=4)


def slide_dashboard(prs, page_num: int, total: int, dashboard: Path):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    content_chrome(s, page_num, total)
    add_slide_header(s, "Daily dashboard (real)",
                     "The first thing you see when you open the app.")
    add_bullets(s, Inches(0.6), Inches(2.0), Inches(7.5), Inches(5.0), [
        ("Readiness score at the top. ", "Color-coded band tells you what kind of day it is in one glance."),
        ("Week strip. ", "Today highlighted; tap any day to jump."),
        ("Key metrics in a grid. ", "HRV, resting HR, sleep, strain, weight, body fat — all with day-over-day deltas."),
        ("Impact donut. ", "Why your readiness is what it is — every input's contribution mapped."),
        ("Trends below the fold. ", "Steps, calories, weight charts auto-stream below."),
        ("FAB for quick log. ", "Drop a mood, hydration, or note without leaving the screen."),
    ], size=13, gap=6)
    add_phone_frame(s, dashboard, Inches(8.7), Inches(1.85), Inches(4.2), Inches(5.3))


def slide_through_day(prs, page_num: int, total: int, log: Path, strength: Path, coach: Path):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    content_chrome(s, page_num, total)
    add_slide_header(s, "Through the day (real screens)",
                     "Strength, food, and the coach — mostly hands-off.")
    phone_w = Inches(3.5)
    phone_h = Inches(5.0)
    positions = [
        (Inches(0.7),  log,      "/log",      "Macros bar live from Yazio.\nWHOOP sleep + nutrition + mood in one daily-log row."),
        (Inches(4.95), strength, "/strength", "Today's plan, set logging, anatomy highlight.\nClick an exercise → see the muscles it works."),
        (Inches(9.2),  coach,    "/coach",    "Today / Recent / Tools tabs · active-block progress.\nReview-ready banner. Brief + chat in one timeline."),
    ]
    for x, img, name, sub in positions:
        add_phone_frame(s, img, x, Inches(1.85), phone_w, phone_h)
        add_text(s, x, Inches(7.0), phone_w, Inches(0.35),
                 name, size=15, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        add_text(s, x, Inches(7.32), phone_w, Inches(0.45),
                 sub, size=10, color=MID, align=PP_ALIGN.CENTER, line_spacing=1.3)


def slide_cadence(prs, page_num: int, total: int, weekly: Path, trends: Path):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    content_chrome(s, page_num, total)
    add_slide_header(s, "Weekly & monthly cadence",
                     "Sunday's review document closes one week and prescribes the next.")
    add_bullets(s, Inches(0.6), Inches(1.95), Inches(7.3), Inches(5.0), [
        ("Weekly Review (Sunday cron). ", "Auto-generates a versioned document per (user, week_start). Recap → Reconfirm → Trends → Prescription → Targets."),
        ("Three surfaces, one doc. ", "Chat card on Sunday · full page at /coach/weeks/[date] · Tue-Sat banner on /coach if you haven't answered yet."),
        ("Commit writes next week. ", "Approval-token commit writes the next training_weeks row and stamps the review with committed_training_week_id."),
        ("Mid-week swaps. ", "One-tap A↔B exchange or single-day replacement. Original Sunday plan preserved (original_session_plan) so adherence math stays anchored."),
        ("Monthly — body measurements. ", "14 circumferences + private photo. Tracks composition drift the scale can't see."),
        ("Every 5 weeks — new block. ", "Pick the next mesocycle goal. Coach respects medical context (GLP-1 phase, injuries)."),
    ], size=13, gap=6)
    # Mockup + real trends
    add_text(s, Inches(8.0), Inches(1.85), Inches(2.5), Inches(0.3),
             "PLAN NEXT WEEK", size=10, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_image_fit(s, weekly, Inches(7.8), Inches(2.2), Inches(2.4), Inches(4.9))
    add_text(s, Inches(10.5), Inches(1.85), Inches(2.5), Inches(0.3),
             "TRENDS (REAL)", size=10, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_phone_frame(s, trends, Inches(10.4), Inches(2.2), Inches(2.5), Inches(4.9))


def slide_coach_speaks(prs, page_num: int, total: int, coach_trends: Path | None):
    """The proactive layer + /coach/trends — the third rhythm."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    content_chrome(s, page_num, total)
    add_slide_header(s, "The coach speaks first",
                     "Deterministic triggers fire, chat cards land, /coach/trends explains why.")

    # Left column — proactive triggers (3 cards)
    add_text(s, Inches(0.6), Inches(1.95), Inches(5.4), Inches(0.3),
             "PROACTIVE NUDGES — DAILY 11:00 UTC CRON",
             size=10, bold=True, color=ACCENT)

    triggers = [
        ("Plateau", "A Big-4 lift's e1RM has been flat past the configured window.", DANGER),
        ("Off-pace weight", "Body-weight slope materially off the block's target (cut / bulk / recomp).", WARNING),
        ("HRV below baseline", "Rolling HRV has dropped vs the WHOOP baseline.", ACCENT),
    ]
    y = Inches(2.35)
    for name, body, color in triggers:
        card = add_rect(s, Inches(0.6), y, Inches(5.4), Inches(1.3),
                        fill=SURFACE, line=(232, 234, 243), radius=0.05, line_width=0.5)
        card.shadow.inherit = False
        add_rect(s, Inches(0.6), y, Inches(0.08), Inches(1.3), fill=color)
        add_pill(s, Inches(0.85), y + Inches(0.22), Inches(1.7), Inches(0.36),
                 name.upper(), fill=color, text_color=(255, 255, 255), size=10)
        add_text(s, Inches(0.85), y + Inches(0.68), Inches(4.4), Inches(0.55),
                 body, size=11, color=MID, line_spacing=1.3)
        y += Inches(1.45)

    # Middle column — /coach/trends sections
    add_text(s, Inches(6.4), Inches(1.95), Inches(3.7), Inches(0.3),
             "/COACH/TRENDS — DEEP COACHING SURFACE",
             size=10, bold=True, color=ACCENT)

    sections = [
        ("Performance",
         "Per-lift e1RM with OLS slopes, plateau spans, strength-per-LBM, IPF GL."),
        ("Composition",
         "Weight, lean mass, fat mass vs the block's targets — off-pace flags fire here."),
        ("Cross",
         "Plain-English insights linking metrics. Deterministic templating, no AI calls."),
    ]
    y = Inches(2.35)
    for name, body in sections:
        card = add_rect(s, Inches(6.4), y, Inches(3.7), Inches(1.3),
                        fill=SURFACE, line=(232, 234, 243), radius=0.05, line_width=0.5)
        card.shadow.inherit = False
        add_rect(s, Inches(6.4), y, Inches(0.08), Inches(1.3), fill=ACCENT)
        add_text(s, Inches(6.65), y + Inches(0.22), Inches(3.3), Inches(0.4),
                 name, size=14, bold=True, color=STRONG)
        add_text(s, Inches(6.65), y + Inches(0.62), Inches(3.3), Inches(0.7),
                 body, size=10, color=MID, line_spacing=1.3)
        y += Inches(1.45)

    # Right column — real /coach/trends screenshot
    if coach_trends is not None:
        add_text(s, Inches(10.5), Inches(1.85), Inches(2.4), Inches(0.3),
                 "/COACH/TRENDS (REAL)", size=10, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        add_phone_frame(s, coach_trends, Inches(10.4), Inches(2.2), Inches(2.5), Inches(4.9))

    # Footer band — cards deep-link
    add_text(s, Inches(0.6), Inches(6.85), Inches(12.1), Inches(0.4),
             "Chat cards deep-link into the matching /coach/trends section. "
             "Dedup is the chat history itself — 7-day window, audit trail built in.",
             size=12, italic=True, color=MID, align=PP_ALIGN.CENTER)


def slide_mental_model(prs, page_num: int, total: int):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    content_chrome(s, page_num, total)
    add_slide_header(s, "The mental model",
                     "A coach that already knows your file — and watches the streams.")

    sources = [
        ("WHOOP",        (20, 184, 112)),
        ("Withings",     (0, 175, 255)),
        ("Apple Health", (15, 20, 48)),
        ("Strong",       (239, 68, 68)),
        ("Yazio",        (249, 115, 22)),
    ]
    top_y = Inches(2.3)
    src_x = Inches(0.6)
    src_w = Inches(2.4)
    src_h = Inches(0.65)
    for i, (name, color) in enumerate(sources):
        y = top_y + Inches(i * 0.85)
        add_pill(s, src_x, y, src_w, src_h, name, fill=color, text_color=(255, 255, 255), size=13)

    # Center brain
    brain_x = Inches(4.4)
    brain_y = Inches(2.6)
    brain_w = Inches(4.5)
    brain_h = Inches(2.5)
    box = add_rect(s, brain_x, brain_y, brain_w, brain_h, fill=SURFACE, line=(232, 234, 243), radius=0.04)
    box.shadow.inherit = False
    add_rect(s, brain_x, brain_y, Inches(0.08), brain_h, fill=ACCENT)
    add_text(s, brain_x + Inches(0.3), brain_y + Inches(0.3), brain_w - Inches(0.6), Inches(0.4),
             "APEX HEALTH OS", size=11, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_text(s, brain_x + Inches(0.3), brain_y + Inches(0.7), brain_w - Inches(0.6), Inches(0.6),
             "Acknowledged profile", size=20, bold=True, color=STRONG, align=PP_ALIGN.CENTER)
    add_text(s, brain_x + Inches(0.3), brain_y + Inches(1.25), brain_w - Inches(0.6), Inches(0.6),
             "+ active plan", size=20, bold=True, color=STRONG, align=PP_ALIGN.CENTER)
    add_text(s, brain_x + Inches(0.3), brain_y + Inches(1.9), brain_w - Inches(0.6), Inches(0.4),
             "Deterministic composer · Claude-narrated coach",
             size=11, color=MUTED, align=PP_ALIGN.CENTER)

    # Right column: touchpoints — three rhythms
    tp_x = Inches(9.6)
    add_pill(s, tp_x, Inches(2.25), Inches(3.3), Inches(0.6),
             "☀  Morning brief", fill=(254, 243, 199), text_color=(146, 64, 14), size=13)
    add_text(s, tp_x, Inches(2.92), Inches(3.3), Inches(0.4),
             "Sets the day.", size=10, color=MID, align=PP_ALIGN.CENTER)
    add_pill(s, tp_x, Inches(3.55), Inches(3.3), Inches(0.6),
             "⚡  Proactive nudge", fill=ACCENT_SOFT, text_color=ACCENT_DEEP, size=13)
    add_text(s, tp_x, Inches(4.22), Inches(3.3), Inches(0.4),
             "Pings when a trigger fires.", size=10, color=MID, align=PP_ALIGN.CENTER)
    add_pill(s, tp_x, Inches(4.85), Inches(3.3), Inches(0.6),
             "☾  Sleep capture", fill=(46, 58, 140), text_color=(255, 255, 255), size=13)
    add_text(s, tp_x, Inches(5.52), Inches(3.3), Inches(0.4),
             "Closes the loop overnight.", size=10, color=MID, align=PP_ALIGN.CENTER)

    # Connector lines from sources → brain
    for i in range(len(sources)):
        y = top_y + Inches(i * 0.85) + Inches(0.325)
        connector = s.shapes.add_connector(1, src_x + src_w, y, brain_x, brain_y + Inches(1.25))
        connector.line.color.rgb = rgb((180, 184, 200))
        connector.line.width = Pt(1)
    # Brain → touchpoints (3 rhythms)
    for ty in (Inches(2.55), Inches(3.85), Inches(5.15)):
        connector = s.shapes.add_connector(1, brain_x + brain_w, brain_y + Inches(1.25), tp_x, ty)
        connector.line.color.rgb = rgb((180, 184, 200))
        connector.line.width = Pt(1)

    add_text(s, Inches(0.6), Inches(6.45), Inches(12.1), Inches(0.45),
             "Everything else removes friction from data entry — the human keeps the calls",
             size=14, bold=True, color=STRONG, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.6), Inches(6.85), Inches(12.1), Inches(0.4),
             "(commit week · commit block · acknowledge milestones).",
             size=12, italic=True, color=MID, align=PP_ALIGN.CENTER)


# ─── main ─────────────────────────────────────────────────────────────────────

def build():
    if not Path(CHROME).exists():
        raise SystemExit(f"Chromium not found at {CHROME}")
    print("Rendering mockup scenes…")
    mockups = render_all_mockups()

    print("Loading real screenshots…")
    real = {}
    for name in ["dashboard", "log", "strength", "trends", "coach", "coach-trends", "profile"]:
        p = REAL / f"{name}.png"
        if not p.exists():
            print(f"  WARNING: missing {p} — re-run capture-real.mjs")
        else:
            real[name] = p
            print(f"  found {name}.png")

    print("\nAssembling presentation…")
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    total = 10  # numbered slides after the cover
    slide_title(prs, mockups["brief"])
    slide_why(prs, 1, total)
    slide_summary(prs, 2, total)
    slide_first_open(prs, 3, total, mockups["onboarding"])
    slide_connect(prs, 4, total, real.get("profile", mockups["integrations"]), mockups["integrations"])
    slide_morning(prs, 5, total, mockups["intake"], mockups["brief"])
    slide_dashboard(prs, 6, total, real.get("dashboard", mockups["brief"]))
    slide_through_day(prs, 7, total,
                      real.get("log", mockups["brief"]),
                      real.get("strength", mockups["brief"]),
                      real.get("coach", mockups["brief"]))
    slide_cadence(prs, 8, total, mockups["weekly"], real.get("trends", mockups["brief"]))
    slide_coach_speaks(prs, 9, total, real.get("coach-trends"))
    slide_mental_model(prs, 10, total)

    prs.save(OUT)
    print(f"\nWrote {OUT}")


if __name__ == "__main__":
    build()
